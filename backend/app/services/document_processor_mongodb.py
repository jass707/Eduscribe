"""
Document processing service with MongoDB vector storage
Replaces FAISS with MongoDB Atlas Vector Search
"""
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import numpy as np

from app.core.config import settings
from database.mongodb_connection import (
    save_document,
    save_document_embeddings,
    vector_search,
    simple_vector_search,
    mark_document_processed
)

# Global embedder (lazy loaded)
_embedder = None

def get_embedder():
    """Get or create the sentence transformer model."""
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer(settings.EMBEDDING_MODEL)
    return _embedder

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from PDF file."""
    try:
        reader = PdfReader(pdf_path)
        text = []
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text.append(content)
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting PDF {pdf_path}: {e}")
        return ""

def extract_text_from_ppt(ppt_path: str) -> str:
    """Extract text from PPTX file."""
    try:
        prs = Presentation(ppt_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text.append(shape.text.strip())
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting PPT {ppt_path}: {e}")
        return ""

def extract_text_from_docx(docx_path: str) -> str:
    """Extract text from DOCX file."""
    try:
        doc = docx.Document(docx_path)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text.strip())
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting DOCX {docx_path}: {e}")
        return ""

def extract_text_from_txt(txt_path: str) -> str:
    """Extract text from TXT file."""
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"Error extracting TXT {txt_path}: {e}")
        return ""

def chunk_text(text: str, chunk_size: int = 300) -> List[str]:
    """Split text into word chunks."""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk.strip())
    return chunks

async def process_document(
    file_path: str,
    lecture_id: str,
    filename: str
) -> Dict[str, Any]:
    """
    Process a document and store in MongoDB with embeddings.
    
    Args:
        file_path: Path to the document file
        lecture_id: ID of the lecture this document belongs to
        filename: Original filename
    
    Returns:
        Dict with document_id, chunk_count, and status
    """
    print(f"📄 Processing document: {filename}")
    
    # Extract text based on file type
    file_ext = Path(file_path).suffix.lower()
    
    if file_ext == '.pdf':
        text = extract_text_from_pdf(file_path)
        file_type = 'pdf'
    elif file_ext in ['.ppt', '.pptx']:
        text = extract_text_from_ppt(file_path)
        file_type = 'pptx'
    elif file_ext in ['.doc', '.docx']:
        text = extract_text_from_docx(file_path)
        file_type = 'docx'
    elif file_ext == '.txt':
        text = extract_text_from_txt(file_path)
        file_type = 'txt'
    else:
        return {
            "success": False,
            "error": f"Unsupported file type: {file_ext}"
        }
    
    if not text or len(text.strip()) < 50:
        return {
            "success": False,
            "error": "No text extracted or text too short"
        }
    
    print(f"✅ Extracted {len(text)} characters from {filename}")
    
    # Save document metadata to MongoDB
    document_id = await save_document(
        lecture_id=lecture_id,
        filename=filename,
        file_type=file_type,
        file_path=file_path,
        content=text
    )
    
    print(f"✅ Saved document to MongoDB: {document_id}")
    
    # Chunk the text
    chunks = chunk_text(text, chunk_size=300)
    print(f"✅ Created {len(chunks)} chunks")
    
    # Generate embeddings
    embedder = get_embedder()
    embeddings = embedder.encode(chunks, show_progress_bar=False)
    print(f"✅ Generated embeddings: {embeddings.shape}")
    
    # Prepare data for MongoDB
    embedding_data = [
        {
            'lecture_id': lecture_id,
            'document_id': document_id,
            'chunk_text': chunk,
            'chunk_index': i,
            'embedding': embedding,
            'metadata': {
                'filename': filename,
                'file_type': file_type
            }
        }
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings))
    ]
    
    # Save embeddings to MongoDB
    await save_document_embeddings(embedding_data)
    print(f"✅ Saved {len(chunks)} embeddings to MongoDB")
    
    # Mark document as processed
    await mark_document_processed(document_id)
    
    return {
        "success": True,
        "document_id": document_id,
        "chunk_count": len(chunks),
        "text_length": len(text)
    }

async def query_documents(
    query_text: str,
    lecture_id: str,
    top_k: int = 10,
    use_atlas_search: bool = True
) -> List[str]:
    """
    Query documents using vector similarity search.
    
    Args:
        query_text: The query text (transcription)
        lecture_id: ID of the lecture
        top_k: Number of top results to return
        use_atlas_search: Try Atlas Vector Search first, fallback to simple search
    
    Returns:
        List of relevant text chunks
    """
    # Generate query embedding
    embedder = get_embedder()
    query_embedding = embedder.encode(query_text, show_progress_bar=False)
    
    # Try Atlas Vector Search first
    if use_atlas_search:
        try:
            results = await vector_search(
                query_embedding=query_embedding,
                lecture_id=lecture_id,
                top_k=top_k
            )
            print(f"✅ Atlas Vector Search returned {len(results)} results")
            return [r['chunk_text'] for r in results]
        except Exception as e:
            print(f"⚠️  Atlas Vector Search failed, using fallback: {e}")
    
    # Fallback to simple cosine similarity
    results = await simple_vector_search(
        query_embedding=query_embedding,
        lecture_id=lecture_id,
        top_k=top_k
    )
    
    print(f"✅ Simple vector search returned {len(results)} results")
    return [r['chunk_text'] for r in results]

# Backward compatibility: Keep the old function name
async def query_documents_faiss(query_text: str, lecture_id: str, top_k: int = 10) -> List[str]:
    """Backward compatibility wrapper for query_documents"""
    return await query_documents(query_text, lecture_id, top_k)
